

from .PPTX_Slide import PPTX_Slide, PPTX_TextLine
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.action import PP_ACTION_TYPE
from .PPTX_Utils import *
from pathlib import Path
import json


class PPTX_Generator:

    SLIDE_LAYOUT_ONLY_TITLE = 5
    SLIDE_LAYOUT_EMPTY = 6

    def __init__(self, p_size: str = '16:10'):

        self.h_presentation = Presentation()

        if p_size == '16:10':
            w, h = 10, 10*10/16
        elif p_size == '16:9':
            w, h = 10, 10 * 9 / 16
        elif p_size == '4:3':
            w, h = 10, 10 * 3 / 4
        else:
            assert False, f'p_size must be one of the following option: ["16:10",  "16:9" , "4:3"] !!!'

        self.h_presentation.slide_width = Inches(w)
        self.h_presentation.slide_height = Inches(h)
        self.size = (self.h_presentation.slide_width, self.h_presentation.slide_height)
        self.slide_layout = self.h_presentation.slide_layouts[self.SLIDE_LAYOUT_EMPTY]


    def create_slide_template(self, slide_template_filename):

        with open(slide_template_filename) as jf:
            configuration = json.load(jf)

        h_slide = PPTX_Slide(self.size, configuration)
        return h_slide

    def add_slide(self, slide: PPTX_Slide):
        new_slide = self.h_presentation.slides.add_slide(self.slide_layout)
        self._add_top_lines(new_slide, slide.top_lines)
        self._add_images(new_slide, slide.images)

    def _get_alignment(self, alignment_str: str):
        if alignment_str == 'center':
            val = PP_PARAGRAPH_ALIGNMENT.CENTER
        elif alignment_str == 'left':
            val = PP_PARAGRAPH_ALIGNMENT.LEFT
        elif alignment_str == 'right':
            val = PP_PARAGRAPH_ALIGNMENT.RIGHT
        else:
            assert 0, f'Unknown alignment {alignment_str}.'
        return val

    def _add_top_lines(self, new_slide, top_lines):
        for line in top_lines:
            textbox = new_slide.shapes.add_textbox(*line.ltwh)
            paragraph = textbox.text_frame.paragraphs[0]

            hlink = line.hlink
            if hlink:
                r = paragraph.add_run()
                r.text = line.text
                r.hyperlink.address = hlink
            else:
                textbox.text = line.text
                paragraph = textbox.text_frame.paragraphs[0]

            tb_font = paragraph.font

            tb_font.name = line.font['name']
            tb_font.size = Pt(line.font['size'])
            tb_font.underline = (line.font['underline'] > 0)

            paragraph.alignment = self._get_alignment(line.font['alignment'])

    def _add_images(self, new_slide, images):
        for images_y in images:
            for image in images_y:
                image_path = image.file_path
                if len(image_path)>0:
                    image_ph = new_slide.shapes.add_picture(image_path, *image.ltwh)
                    action = image.hlink
                    if action:
                        image_ph.click_action.hyperlink.address = action # image_ph.click_action.action = PP_ACTION_TYPE.RUN_PROGRAM

    def save(self, filepath):
        filepath = Path(filepath)
        temp_path = create_new_temp_directory_and_remove_old_dirs()
        temp_file_path = temp_path / filepath.name
        self.h_presentation.save(temp_file_path.as_posix())
        unzip_file(temp_file_path)
        update_action_to_program(temp_path)
        zip_folder(temp_path, filepath)
        delete_path(temp_path)


if __name__=="__main__":
    image_paths = [ r'C:\Users\Public\Pictures\Sample Pictures\Hydrangeas.jpg',
                    r'C:\Users\Public\Pictures\Sample Pictures\Chrysanthemum.jpg',
                   r'C:\Users\Public\Pictures\Sample Pictures\Desert.jpg',
                    r'C:\Users\Public\Pictures\Sample Pictures\Penguins.jpg']


    configuration_filename = "Slide_1_Template.json"
    h_pptx = PPTX_Generator(configuration_filename)


    h_pptx.slide.set_line(0, 'Main Title')
    h_pptx.slide.set_line(1, 'Second line...')
    h_pptx.slide.set_line(2, 'Last Line !!!')


    actions = ['word', 'cmd', 'calc', 'd:/Kill_Excel.bat']

    for y in range(4):
        for x in range(3):
            h_pptx.slide.set_image(x, y, image_paths[(x+y*3)%4])

    h_pptx.add_slide()

    h_pptx.save('P1.pptx')
